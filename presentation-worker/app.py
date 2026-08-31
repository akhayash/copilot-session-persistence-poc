from __future__ import annotations

import hashlib
import base64
import io
import json
import os
import re
import signal
import shutil
import subprocess
import sys
import tempfile
import threading
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

import fitz
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import FileResponse, JSONResponse
from PIL import Image
from pydantic import BaseModel, ConfigDict, Field, field_validator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

MAX_REQUEST_BYTES = 64 * 1024
MAX_ARTIFACT_BYTES = 32 * 1024 * 1024
MAX_EXEC_OUTPUT_BYTES = 64 * 1024
MAX_RENDER_IMAGES = 12
MAX_RENDER_BYTES = 8 * 1024 * 1024
MAX_PREVIEW_EDGE = 1024
SAFE_FILE_NAME = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,119}\.pptx$", re.IGNORECASE)
SLIDE_XML = re.compile(r"^ppt/slides/slide\d+\.xml$")
INTEGER_COORDINATE = re.compile(r"^-?\d+$")
CONTENT_TYPES = {
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".json": "application/json",
}

NAVY = RGBColor(24, 39, 66)
INK = RGBColor(36, 45, 58)
CREAM = RGBColor(247, 244, 236)
TEAL = RGBColor(24, 137, 141)
CORAL = RGBColor(229, 107, 87)
PALE_TEAL = RGBColor(221, 239, 237)
WHITE = RGBColor(255, 255, 255)
FONT = "Noto Sans CJK JP"

WORKSPACE_DIR = Path(
    os.environ.get("PRESENTATION_WORKSPACE", tempfile.mkdtemp(prefix="presentation-worker-"))
).resolve()
# Kept as an alias so older tests and callers can override the legacy artifact directory.
OUTPUT_DIR = WORKSPACE_DIR
BUILD_LOCK = threading.Lock()


class ContentSlide(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str = Field(min_length=1, max_length=100)
    body: str = Field(min_length=1, max_length=800)
    highlight: str | None = Field(default=None, min_length=1, max_length=180)

    @field_validator("title", "body", "highlight")
    @classmethod
    def reject_blank(cls, value: str | None) -> str | None:
        if value is not None and not value.strip():
            raise ValueError("must not be blank")
        return value


class PresentationRequest(BaseModel):
    model_config = ConfigDict(extra="forbid")

    fileName: str = Field(min_length=6, max_length=125)
    title: str = Field(min_length=1, max_length=140)
    subtitle: str | None = Field(default=None, min_length=1, max_length=240)
    audience: str = Field(min_length=1, max_length=160)
    slides: list[ContentSlide] = Field(min_length=1, max_length=7)

    @field_validator("fileName")
    @classmethod
    def safe_file_name(cls, value: str) -> str:
        if not SAFE_FILE_NAME.fullmatch(value) or Path(value).name != value:
            raise ValueError("must be a safe .pptx base file name")
        return value

    @field_validator("title", "subtitle", "audience")
    @classmethod
    def reject_blank(cls, value: str | None) -> str | None:
        if value is not None and not value.strip():
            raise ValueError("must not be blank")
        return value


class ExecRequest(BaseModel):
    model_config = ConfigDict(extra="forbid")

    command: str = Field(min_length=1, max_length=16_384)
    timeoutSeconds: int = Field(default=60, ge=1, le=90)

    @field_validator("command")
    @classmethod
    def reject_blank_command(cls, value: str) -> str:
        if not value.strip():
            raise ValueError("must not be blank")
        return value


class FileWriteRequest(BaseModel):
    model_config = ConfigDict(extra="forbid")

    data: str = Field(max_length=48 * 1024 * 1024)
    encoding: str = Field(default="base64", pattern="^(base64|utf-8)$")


class RenderRequest(BaseModel):
    model_config = ConfigDict(extra="forbid")

    path: str = Field(min_length=6, max_length=240)


app = FastAPI(title="Presentation Worker", docs_url=None, redoc_url=None)


@app.middleware("http")
async def bound_request(request: Request, call_next):
    length = request.headers.get("content-length")
    if request.method == "POST" and request.url.path == "/presentations" and length is None:
        return JSONResponse(status_code=411, content={"detail": "Content-Length is required"})
    if length is not None:
        try:
            maximum = (
                MAX_ARTIFACT_BYTES * 2
                if request.method == "PUT" and request.url.path.startswith("/files/")
                else MAX_REQUEST_BYTES
            )
            if int(length) > maximum:
                return JSONResponse(status_code=413, content={"detail": "Request body is too large"})
        except ValueError:
            return JSONResponse(status_code=400, content={"detail": "Invalid Content-Length"})
    return await call_next(request)


@app.get("/healthz")
def healthz() -> dict[str, str]:
    return {"status": "ok"}


def workspace_root() -> Path:
    return OUTPUT_DIR.resolve()


def workspace_path(relative_path: str, *, must_exist: bool = False) -> Path:
    if not relative_path or "\x00" in relative_path:
        raise HTTPException(status_code=404, detail="File not found")
    normalized = relative_path.replace("\\", "/")
    if normalized.startswith("/") or any(part in {"", ".", ".."} for part in normalized.split("/")):
        raise HTTPException(status_code=404, detail="File not found")
    root = workspace_root()
    path = (root / normalized).resolve()
    if path == root or root not in path.parents:
        raise HTTPException(status_code=404, detail="File not found")
    if must_exist and not path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return path


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def file_entry(path: Path) -> dict[str, object]:
    return {
        "path": path.relative_to(workspace_root()).as_posix(),
        "sizeBytes": path.stat().st_size,
        "sha256": sha256_file(path),
    }


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    margin: float = 0.06,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_shape(slide, kind, left, top, width, height, fill, line=None, radius=None):
    shape = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def body_parts(body: str) -> list[str]:
    parts = [line.strip() for line in body.splitlines() if line.strip()]
    return parts or [body]


def display_text(value: str) -> str:
    return re.sub(r"([、。])\s+", r"\1", value.strip())


def create_pptx(data: PresentationRequest, destination: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    background = title_slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY
    add_shape(title_slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 0.28, 7.5, CORAL)
    add_shape(title_slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.75, 0.65, 1.45, 1.45, TEAL)
    add_shape(title_slide, MSO_AUTO_SHAPE_TYPE.OVAL, 11.75, 1.55, 0.62, 0.62, CORAL)
    add_text(title_slide, "POWERPOINT ARTIFACT", 0.9, 1.15, 4.0, 0.35, size=12, color=CORAL, bold=True)
    add_text(title_slide, display_text(data.title), 0.85, 2.0, 10.1, 1.05, size=42, color=WHITE, bold=True)
    if data.subtitle:
        add_text(title_slide, display_text(data.subtitle), 0.9, 3.35, 9.5, 0.72, size=20, color=PALE_TEAL)
    add_text(title_slide, f"対象  |  {display_text(data.audience)}", 0.9, 6.35, 8.5, 0.42, size=14, color=WHITE)

    for index, item in enumerate(data.slides, start=1):
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = CREAM
        add_text(slide, f"{index:02d}", 0.65, 0.45, 0.65, 0.38, size=12, color=CORAL, bold=True)
        add_text(slide, display_text(item.title), 1.35, 0.32, 10.9, 0.78, size=32, color=NAVY, bold=True)
        parts = body_parts(item.body)

        if item.highlight and index % 2 == 1:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.75, 1.55, 4.0, 4.65, NAVY)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.1, 1.95, 0.68, 0.68, CORAL)
            add_text(slide, "KEY POINT", 1.95, 2.02, 2.1, 0.42, size=12, color=PALE_TEAL, bold=True)
            add_text(
                slide,
                display_text(item.highlight),
                1.08,
                3.0,
                3.35,
                1.7,
                size=24,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.15, 1.55, 7.35, 4.65, WHITE, PALE_TEAL)
            add_text(slide, display_text(item.body), 5.62, 2.05, 6.4, 3.6, size=19, color=INK)
        elif item.highlight:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.75, 1.5, 11.75, 2.25, WHITE, PALE_TEAL)
            add_text(slide, display_text(item.body), 1.15, 1.88, 10.95, 1.48, size=18, color=INK)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.35, 4.35, 0.72, 0.72, TEAL)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 2.05, 4.68, 1.15, 0.06, TEAL)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.18, 4.05, 6.95, 1.35, NAVY)
            add_text(
                slide,
                display_text(item.highlight),
                3.52,
                4.3,
                6.3,
                0.82,
                size=23,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 10.12, 4.68, 1.15, 0.06, TEAL)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 11.25, 4.35, 0.72, 0.72, CORAL)
        elif index % 2 == 0 and len(parts) > 1:
            count = min(len(parts), 6)
            for part_index, part in enumerate(parts[:count]):
                x = 0.9 + (part_index % 3) * 4.08
                y = 1.75 + (part_index // 3) * 2.25
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.65, 1.75, WHITE, PALE_TEAL)
                add_text(slide, f"{part_index + 1:02d}", x + 0.18, y + 0.14, 0.5, 0.32, size=10, color=CORAL, bold=True)
                add_text(slide, part, x + 0.2, y + 0.5, 3.25, 1.05, size=15, color=INK)
            if len(parts) > count:
                add_text(slide, "\n".join(parts[count:]), 0.9, 6.32, 11.6, 0.55, size=11, color=INK)
        elif len(parts) > 1:
            count = min(len(parts), 7)
            start_y = 1.7
            step = min(0.72, 4.7 / count)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.23, start_y + 0.2, 0.04, step * (count - 1), TEAL)
            for part_index, part in enumerate(parts[:count]):
                y = start_y + part_index * step
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.08, y + 0.13, 0.34, 0.34, CORAL)
                add_text(slide, part, 1.72, y, 10.2, 0.58, size=17, color=INK)
            if len(parts) > count:
                add_text(slide, "\n".join(parts[count:]), 1.72, 6.35, 10.2, 0.42, size=10, color=INK)
        else:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 1.7, 11.55, 4.75, WHITE, PALE_TEAL)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.9, 1.7, 0.16, 4.75, TEAL)
            add_text(slide, item.body, 1.45, 2.05, 10.25, 4.0, size=22, color=INK)

        add_text(
            slide,
            display_text(data.audience),
            8.7,
            6.7,
            3.75,
            0.28,
            size=10,
            color=NAVY,
            align=PP_ALIGN.RIGHT,
        )

    prs.save(destination)


def validate_pptx(path: Path, expected_slides: int) -> dict[str, object]:
    if not path.is_file() or path.stat().st_size == 0:
        raise ValueError("PPTX was not created")
    if path.stat().st_size > MAX_ARTIFACT_BYTES:
        raise ValueError("PPTX exceeds the artifact size limit")
    try:
        with zipfile.ZipFile(path) as archive:
            bad_file = archive.testzip()
            if bad_file:
                raise ValueError(f"Open XML archive contains a corrupt member: {bad_file}")
            names = archive.namelist()
            if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
                raise ValueError("Required Open XML members are missing")
            slides = sorted(name for name in names if SLIDE_XML.fullmatch(name))
            if len(slides) != expected_slides:
                raise ValueError(f"Expected {expected_slides} slides, found {len(slides)}")
            if any(not archive.read(name).strip() for name in slides):
                raise ValueError("A slide XML member is empty")
            for name in slides:
                try:
                    root = ET.fromstring(archive.read(name))
                except ET.ParseError as exc:
                    raise ValueError(f"Slide XML is malformed: {name}") from exc
                for element in root.iter():
                    local_name = element.tag.rsplit("}", 1)[-1]
                    if local_name not in {"off", "ext", "chOff", "chExt"}:
                        continue
                    for attribute in ("x", "y", "cx", "cy"):
                        value = element.attrib.get(attribute)
                        if value is not None and not INTEGER_COORDINATE.fullmatch(value):
                            raise ValueError(
                                "Open XML geometry coordinates must be integers: "
                                f"{name} {local_name}.{attribute}={value!r}"
                            )
    except zipfile.BadZipFile as exc:
        raise ValueError("PPTX is not a valid Open XML zip") from exc
    return {"passed": True, "openXml": True, "nonemptySlideXml": True}


def render_presentation(pptx_path: Path, output_dir: Path, expected_pages: int) -> tuple[Path, list[Path]]:
    profile = output_dir / "lo-profile"
    profile.mkdir()
    command = [
        "libreoffice",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile.resolve().as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(pptx_path),
    ]
    try:
        result = subprocess.run(command, capture_output=True, text=True, timeout=90, check=False)
    except (OSError, subprocess.TimeoutExpired) as exc:
        raise RuntimeError(f"LibreOffice rendering failed: {exc}") from exc
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if result.returncode != 0 or not pdf_path.is_file() or pdf_path.stat().st_size == 0:
        detail = (result.stderr or result.stdout).strip()[-500:]
        raise RuntimeError(f"LibreOffice rendering failed: {detail or 'no PDF produced'}")

    png_paths: list[Path] = []
    try:
        with fitz.open(pdf_path) as document:
            if document.page_count != expected_pages:
                raise RuntimeError(
                    f"Rendered page count differs: expected {expected_pages}, found {document.page_count}"
                )
            for page_number, page in enumerate(document, start=1):
                pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                png_path = output_dir / f"{pptx_path.stem}-slide-{page_number:02d}.png"
                pixmap.save(png_path)
                png_paths.append(png_path)
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError(f"PDF-to-PNG rendering failed: {exc}") from exc

    for artifact in [pdf_path, *png_paths]:
        if not artifact.is_file() or artifact.stat().st_size == 0:
            raise RuntimeError(f"Rendered artifact is empty: {artifact.name}")
        if artifact.stat().st_size > MAX_ARTIFACT_BYTES:
            raise RuntimeError(f"Rendered artifact exceeds size limit: {artifact.name}")
    return pdf_path, png_paths


def artifact_manifest(path: Path) -> dict[str, object]:
    return {
        "fileName": path.name,
        "contentType": CONTENT_TYPES[path.suffix.lower()],
        "sizeBytes": path.stat().st_size,
        "sha256": sha256_file(path),
    }


def clear_outputs() -> None:
    for child in OUTPUT_DIR.iterdir():
        if child.is_dir():
            shutil.rmtree(child)
        else:
            child.unlink()


@app.post("/exec")
def execute(data: ExecRequest) -> dict[str, object]:
    root = workspace_root()
    root.mkdir(parents=True, exist_ok=True)
    try:
        with BUILD_LOCK:
            shell = (
                [os.environ.get("COMSPEC", "cmd.exe"), "/d", "/s", "/c", data.command]
                if os.name == "nt"
                else ["/bin/sh", "-lc", data.command]
            )
            process = subprocess.Popen(
                shell,
                cwd=root,
                stdin=subprocess.DEVNULL,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                start_new_session=os.name != "nt",
            )
            stdout_buffer = bytearray()
            stderr_buffer = bytearray()
            stdout_total = [0]
            stderr_total = [0]
            stdout_thread = threading.Thread(
                target=collect_output,
                args=(process.stdout, stdout_buffer, stdout_total),
                daemon=True,
            )
            stderr_thread = threading.Thread(
                target=collect_output,
                args=(process.stderr, stderr_buffer, stderr_total),
                daemon=True,
            )
            stdout_thread.start()
            stderr_thread.start()
            try:
                exit_code = process.wait(timeout=data.timeoutSeconds)
            except subprocess.TimeoutExpired as exc:
                if os.name == "nt":
                    process.kill()
                else:
                    os.killpg(process.pid, signal.SIGKILL)
                process.wait()
                stdout_thread.join()
                stderr_thread.join()
                raise HTTPException(status_code=408, detail="Command timed out") from exc
            if os.name != "nt":
                try:
                    os.killpg(process.pid, signal.SIGKILL)
                except ProcessLookupError:
                    pass
            stdout_thread.join(timeout=2)
            stderr_thread.join(timeout=2)
            if stdout_thread.is_alive() and process.stdout is not None:
                process.stdout.close()
            if stderr_thread.is_alive() and process.stderr is not None:
                process.stderr.close()
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Command execution failed: {exc}") from exc

    stdout = bytes(stdout_buffer).decode("utf-8", errors="replace")
    stderr = bytes(stderr_buffer).decode("utf-8", errors="replace")
    return {
        "exitCode": exit_code,
        "stdout": stdout,
        "stderr": stderr,
        "stdoutTruncated": stdout_total[0] > MAX_EXEC_OUTPUT_BYTES,
        "stderrTruncated": stderr_total[0] > MAX_EXEC_OUTPUT_BYTES,
    }


def collect_output(stream, buffer: bytearray, total: list[int]) -> None:
    if stream is None:
        return
    with stream:
        for chunk in iter(lambda: stream.read(8192), b""):
            total[0] += len(chunk)
            buffer.extend(chunk)
            if len(buffer) > MAX_EXEC_OUTPUT_BYTES:
                del buffer[:-MAX_EXEC_OUTPUT_BYTES]


@app.get("/files")
def list_files() -> dict[str, object]:
    root = workspace_root()
    root.mkdir(parents=True, exist_ok=True)
    files = [
        file_entry(path)
        for path in sorted(root.rglob("*"))
        if path.is_file()
    ]
    return {"files": files}


@app.get("/files/{relative_path:path}")
def read_file(relative_path: str):
    path = workspace_path(relative_path, must_exist=True)
    if path.stat().st_size > MAX_ARTIFACT_BYTES:
        raise HTTPException(status_code=413, detail="File exceeds the download size limit")
    return FileResponse(path, media_type="application/octet-stream", filename=path.name)


@app.put("/files/{relative_path:path}")
def write_file(relative_path: str, data: FileWriteRequest) -> dict[str, object]:
    path = workspace_path(relative_path)
    try:
        content = (
            base64.b64decode(data.data, validate=True)
            if data.encoding == "base64"
            else data.data.encode("utf-8")
        )
    except (ValueError, TypeError) as exc:
        raise HTTPException(status_code=422, detail="Invalid base64 file content") from exc
    if len(content) > MAX_ARTIFACT_BYTES:
        raise HTTPException(status_code=413, detail="File exceeds the upload size limit")
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            dir=path.parent,
            prefix=".upload-",
            delete=False,
        ) as temporary:
            temporary.write(content)
            temporary_path = Path(temporary.name)
        temporary_path.replace(path)
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)
    return file_entry(path)


@app.delete("/files/{relative_path:path}", status_code=204)
def delete_file(relative_path: str):
    path = workspace_path(relative_path, must_exist=True)
    path.unlink()
    parent = path.parent
    root = workspace_root()
    while parent != root and not any(parent.iterdir()):
        parent.rmdir()
        parent = parent.parent


def preview_image(path: Path) -> tuple[str, int]:
    with Image.open(path) as image:
        image.thumbnail((MAX_PREVIEW_EDGE, MAX_PREVIEW_EDGE), Image.Resampling.LANCZOS)
        output = io.BytesIO()
        image.convert("RGB").save(output, format="PNG", optimize=True)
    content = output.getvalue()
    return base64.b64encode(content).decode("ascii"), len(content)


@app.post("/render")
def render(data: RenderRequest) -> dict[str, object]:
    pptx_path = workspace_path(data.path, must_exist=True)
    if pptx_path.suffix.lower() != ".pptx":
        raise HTTPException(status_code=422, detail="Only .pptx files can be rendered")
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            expected_pages = sum(1 for name in archive.namelist() if SLIDE_XML.fullmatch(name))
    except zipfile.BadZipFile as exc:
        raise HTTPException(status_code=422, detail="PPTX is not a valid Open XML zip") from exc
    if expected_pages < 1 or expected_pages > MAX_RENDER_IMAGES:
        raise HTTPException(status_code=422, detail="Presentation slide count is outside render limits")

    render_dir = Path(tempfile.mkdtemp(prefix="presentation-render-")).resolve()
    try:
        with BUILD_LOCK:
            validation = validate_pptx(pptx_path, expected_pages)
            pdf_path, png_paths = render_presentation(pptx_path, render_dir, expected_pages)
        images = []
        total_bytes = 0
        for index, png_path in enumerate(png_paths, start=1):
            encoded, size_bytes = preview_image(png_path)
            total_bytes += size_bytes
            if total_bytes > MAX_RENDER_BYTES:
                raise HTTPException(status_code=413, detail="Rendered previews exceed the response size limit")
            images.append(
                {
                    "slideNumber": index,
                    "mimeType": "image/png",
                    "data": encoded,
                    "sizeBytes": size_bytes,
                }
            )
        result = {
            "validation": validation,
            "slideCount": expected_pages,
            "pdf": artifact_manifest(pdf_path),
            "images": images,
        }
        return result
    except HTTPException:
        raise
    except (ValueError, RuntimeError) as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    finally:
        shutil.rmtree(render_dir, ignore_errors=True)


@app.post("/presentations")
def presentations(data: PresentationRequest) -> dict[str, object]:
    expected_slides = len(data.slides) + 1
    with BUILD_LOCK:
        workspace_root().mkdir(parents=True, exist_ok=True)
        clear_outputs()
        pptx_path = OUTPUT_DIR / data.fileName
        try:
            create_pptx(data, pptx_path)
            validation = validate_pptx(pptx_path, expected_slides)
            pdf_path, png_paths = render_presentation(pptx_path, OUTPUT_DIR, expected_slides)
            artifacts = [pptx_path, pdf_path, *png_paths]
            artifact_entries = [artifact_manifest(path) for path in artifacts]
            validation_path = OUTPUT_DIR / "validation.json"
            validation_path.write_text(
                json.dumps(
                    {
                        "validationPassed": True,
                        "validation": validation,
                        "slideCount": expected_slides,
                        "files": artifact_entries,
                    },
                    ensure_ascii=False,
                    indent=2,
                )
                + "\n",
                encoding="utf-8",
            )
            return {
                "validationPassed": True,
                "validation": validation,
                "slideCount": expected_slides,
                "files": [*artifact_entries, artifact_manifest(validation_path)],
            }
        except (ValueError, RuntimeError) as exc:
            clear_outputs()
            raise HTTPException(status_code=500, detail=str(exc)) from exc
        except Exception as exc:
            clear_outputs()
            raise HTTPException(status_code=500, detail="Presentation generation failed") from exc


@app.get("/artifacts/{file_name}")
def artifact(file_name: str):
    if Path(file_name).name != file_name or file_name in {".", ".."}:
        raise HTTPException(status_code=404, detail="Artifact not found")
    path = (OUTPUT_DIR / file_name).resolve()
    if path.parent != OUTPUT_DIR or path.suffix.lower() not in CONTENT_TYPES or not path.is_file():
        raise HTTPException(status_code=404, detail="Artifact not found")
    return FileResponse(
        path,
        media_type=CONTENT_TYPES[path.suffix.lower()],
        filename=path.name,
        content_disposition_type="attachment",
    )
